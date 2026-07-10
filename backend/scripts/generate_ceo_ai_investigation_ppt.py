"""Generate 5-slide factual CEO briefing — themes by understanding, not ticket titles.

Run: python -m scripts.generate_ceo_ai_investigation_ppt
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from scripts.ppt_theme_analysis import run as analyze_themes

ROOT = Path(__file__).resolve().parents[2]
OUT_PATH = ROOT / "AI-Investigation-CEO-Briefing.pptx"

BG_DARK = RGBColor(0x1A, 0x1A, 0x1A)
BG_CARD = RGBColor(0x26, 0x26, 0x26)
ACCENT = RGBColor(0xF9, 0x73, 0x16)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA3, 0xA3, 0xA3)
BLUE = RGBColor(0x60, 0xA5, 0xFA)


def set_slide_bg(slide, color: RGBColor = BG_DARK) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, *, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Segoe UI"
    p.alignment = align


def add_bullets(slide, bullets: list[str], top=Inches(1.55), size=15):
    box = slide.shapes.add_textbox(Inches(0.65), top, Inches(8.7), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = WHITE
        p.font.name = "Segoe UI"
        p.space_after = Pt(12)
        p.level = 0


def header(slide, title: str, subtitle: str = ""):
    add_text(slide, Inches(0.6), Inches(0.4), Inches(8.8), Inches(0.65), title, size=30, bold=True)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.05), Inches(8.8), Inches(0.45), subtitle, size=13, color=GRAY)
    bar = slide.shapes.add_shape(1, Inches(0.6), Inches(1.42), Inches(1.0), Inches(0.035))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def add_theme_table(slide, rows: list[dict], highlight: str | None = None):
    headers = ["Product area", "Pre-AI /mo", "Post-AI /mo", "What reports describe"]
    n = len(rows) + 1
    table = slide.shapes.add_table(n, 4, Inches(0.55), Inches(1.55), Inches(8.9), Inches(0.38 * n + 0.2)).table
    table.columns[0].width = Inches(2.4)
    table.columns[1].width = Inches(1.1)
    table.columns[2].width = Inches(1.1)
    table.columns[3].width = Inches(4.3)
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x2D)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = ACCENT
            p.font.name = "Segoe UI"
    for i, row in enumerate(rows, 1):
        is_hi = highlight and row["name"] == highlight
        for j, val in enumerate([row["name"], str(row["pre_per_month"]), str(row["post_per_month"]), row["description"]]):
            cell = table.cell(i, j)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x30, 0x28, 0x20) if is_hi else (BG_CARD if i % 2 else RGBColor(0x22, 0x22, 0x22))
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9 if j < 3 else 8)
                p.font.color.rgb = BLUE if is_hi and j == 0 else WHITE
                p.font.bold = bool(is_hi and j == 0)
                p.font.name = "Segoe UI"


def build(data: dict) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    vol = data["volume"]
    themes = [t for t in data["themes"] if t["name"] != "Other / cross-cutting"][:7]
    cc = next((t for t in data["themes"] if "Customer Concerns" in t["name"]), None)
    fix = data["fix_areas"]

    # ── Slide 1: Title ──
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_text(s, Inches(0.7), Inches(2.3), Inches(8.6), Inches(0.9), "AI Adoption & Engineering Quality", size=38, bold=True)
    add_text(s, Inches(0.7), Inches(3.2), Inches(8.6), Inches(0.5), "Factual data review for leadership", size=20, color=ACCENT)
    add_text(
        s, Inches(0.7), Inches(4.1), Inches(8.6), Inches(1.4),
        "AI-assisted development from 1 May 2026\n"
        "Pre-AI: Dec 2025 – Apr 2026  ·  Post-AI: May – Jun 2026\n"
        "Source: synced Asana tickets, Issue Intelligence clusters\n"
        "Autorox Delivery Intelligence · July 2026",
        size=13, color=GRAY,
    )

    # ── Slide 2: Volume facts ──
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    header(s, "What changed in volume", "Counts by ticket created date — not a quality verdict")
    add_bullets(s, [
        f"Bugs created per month: {vol['bugs_per_month_pre']} (pre-AI) → {vol['bugs_per_month_post']} (post-AI)",
        f"Enhancements created per month: {vol['enh_per_month_pre']} → {vol['enh_per_month_post']}",
        f"Total bugs in window: {vol['bugs_pre_total']} (5 months pre) · {vol['bugs_post_total']} (2 months post)",
        "June 2026: 74 bugs created — highest single month in the comparison window",
        "July 2026 partial month excluded from monthly averages",
        "No development task in Asana is tagged as AI-assisted or human-only",
    ], top=Inches(1.65), size=16)

    # ── Slide 3: Themes by product understanding ──
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    header(
        s,
        "Where defects cluster — by product area",
        "Grouped by what the tickets describe, not individual ticket titles",
    )
    add_theme_table(s, themes, highlight="Customer Concerns (AI module)" if cc else None)
    if cc:
        add_text(
            s, Inches(0.6), Inches(6.55), Inches(8.8), Inches(0.7),
            "Customer Concerns (AI): work orders not reflecting on concern screen, concern data not updating, "
            "Arabic / localisation input rejected, payload sync issues — "
            f"{cc['post_per_month']}/mo post-AI vs {cc['pre_per_month']}/mo pre-AI.",
            size=10, color=GRAY,
        )

    # ── Slide 4: Recurring patterns ──
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    header(s, "Recurring engineering patterns", "Issue Intelligence — same root fix would resolve these tickets")

    cc_fix = fix.get("Customer Concerns (AI module)", {})
    jc_fix = fix.get("Job Card lifecycle", {})
    inv_fix = fix.get("Invoice & billing", {})
    rep_fix = fix.get("Reports & analytics", {})

    bullets = [
        f"{data['recurring_total']} unique recurring clusters identified across synced tickets",
        f"{data.get('recurring_pre_ai_pct', '—')}% of clusters first seen before 1 May 2026",
        f"{data.get('recurring_new_post_ai', '—')} clusters first seen only after 1 May 2026",
        "",
        "Largest recurring areas (by ticket volume in clusters):",
        f"  · Customer Concerns (AI module) — {cc_fix.get('clusters', 0)} clusters, "
        f"{cc_fix.get('tickets', 0)} tickets, {cc_fix.get('trend_increasing', 0)} trending up",
        f"  · Job Card lifecycle — {jc_fix.get('clusters', 0)} clusters, {jc_fix.get('tickets', 0)} tickets",
        f"  · Invoice & billing — {inv_fix.get('clusters', 0)} clusters, {inv_fix.get('tickets', 0)} tickets",
        f"  · Reports & analytics — {rep_fix.get('clusters', 0)} clusters, {rep_fix.get('tickets', 0)} tickets",
        "",
        "Customer Concerns cluster covers: concern-to-workorder sync, update failures, Arabic language handling",
    ]
    add_bullets(s, bullets, top=Inches(1.6), size=14)

    # ── Slide 5: What data contains / does not ──
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    header(s, "What this data contains — and what it does not", "Factual limits of the current measurement")
    add_bullets(s, [
        "Contains: ticket text, dates, modules, Issue Intelligence clusters, pre/post May 2026 comparison",
        "Contains: enhancement requests tracked separately from product bug clusters",
        "Does not contain: which code changes used AI tools",
        "Does not contain: link from a bug to the specific release or commit that introduced it",
        "Does not establish causation between AI adoption and defect patterns — dates are correlated only",
        "",
        "Customer Concerns is an AI-related product module; its defect rate rose from "
        f"{cc['pre_per_month'] if cc else '—'}/mo to {cc['post_per_month'] if cc else '—'}/mo in the comparison window.",
    ], top=Inches(1.6), size=14)

    return prs


def main() -> None:
    data = analyze_themes()
    prs = build(data)
    prs.save(str(OUT_PATH))
    print(f"Saved 5-slide deck: {OUT_PATH}")


if __name__ == "__main__":
    main()
